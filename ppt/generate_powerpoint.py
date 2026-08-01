import os
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE

def create_deck():
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    blank_layout = prs.slide_layouts[6]

    # Color Palette Definitions
    BG_DARK = RGBColor(0x16, 0x16, 0x16)
    CARD_BG = RGBColor(0x26, 0x26, 0x26)
    CARD_BORDER = RGBColor(0x39, 0x39, 0x39)
    IBM_BLUE = RGBColor(0x0F, 0x62, 0xFE)
    IBM_TEAL = RGBColor(0x00, 0xD7, 0xD2)
    TEXT_WHITE = RGBColor(0xFF, 0xFF, 0xFF)
    TEXT_LIGHT = RGBColor(0xC6, 0xC6, 0xC6)
    TEXT_MUTED = RGBColor(0x8D, 0x8D, 0x8D)

    def set_background(slide):
        bg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), prs.slide_width, prs.slide_height)
        bg.fill.solid()
        bg.fill.fore_color.rgb = BG_DARK
        bg.line.fill.background()

        # Bottom accent strip
        strip = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(7.42), prs.slide_width, Inches(0.08))
        strip.fill.solid()
        strip.fill.fore_color.rgb = IBM_BLUE
        strip.line.fill.background()

    def add_header(slide, title_text, badge_text=""):
        # Header title
        tx_box = slide.shapes.add_textbox(Inches(0.6), Inches(0.4), Inches(10), Inches(0.6))
        tf = tx_box.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.text = title_text.upper()
        p.font.size = Pt(22)
        p.font.bold = True
        p.font.color.rgb = TEXT_WHITE

        if badge_text:
            badge_box = slide.shapes.add_textbox(Inches(10.5), Inches(0.4), Inches(2.2), Inches(0.5))
            btf = badge_box.text_frame
            bp = btf.paragraphs[0]
            bp.alignment = PP_ALIGN.RIGHT
            bp.text = f"[{badge_text}]"
            bp.font.size = Pt(12)
            bp.font.bold = True
            bp.font.color.rgb = IBM_TEAL

        # Divider line
        line = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.6), Inches(1.05), Inches(12.133), Inches(0.02))
        line.fill.solid()
        line.fill.fore_color.rgb = CARD_BORDER
        line.line.fill.background()

    def add_card(slide, left, top, width, height, title, items=None, paragraph_text=None, border_color=CARD_BORDER):
        card = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(left), Inches(top), Inches(width), Inches(height))
        card.fill.solid()
        card.fill.fore_color.rgb = CARD_BG
        card.line.color.rgb = border_color
        card.line.width = Pt(1)

        tx_box = slide.shapes.add_textbox(Inches(left + 0.15), Inches(top + 0.15), Inches(width - 0.3), Inches(height - 0.3))
        tf = tx_box.text_frame
        tf.word_wrap = True
        
        # Card Header
        p_head = tf.paragraphs[0]
        p_head.text = title
        p_head.font.size = Pt(16)
        p_head.font.bold = True
        p_head.font.color.rgb = TEXT_WHITE
        p_head.space_after = Pt(8)

        if paragraph_text:
            p_text = tf.add_paragraph()
            p_text.text = paragraph_text
            p_text.font.size = Pt(13)
            p_text.font.color.rgb = TEXT_LIGHT
            p_text.space_after = Pt(8)

        if items:
            for item in items:
                p_item = tf.add_paragraph()
                p_item.text = f"•  {item}"
                p_item.font.size = Pt(12)
                p_item.font.color.rgb = TEXT_LIGHT
                p_item.space_after = Pt(6)

    def set_speaker_notes(slide, notes_text):
        notes_slide = slide.notes_slide
        tf = notes_slide.notes_text_frame
        tf.text = f"SPEAKER SCRIPT:\n{notes_text}"

    # -------------------------------------------------------------
    # SLIDE 1: TITLE SLIDE
    # -------------------------------------------------------------
    slide1 = prs.slides.add_slide(blank_layout)
    set_background(slide1)
    
    t_box = slide1.shapes.add_textbox(Inches(1.0), Inches(1.2), Inches(11.333), Inches(4.5))
    tf1 = t_box.text_frame
    tf1.word_wrap = True

    p0 = tf1.paragraphs[0]
    p0.text = "IBM CONSULTING STAND & DELIVER PRESENTATION"
    p0.font.size = Pt(14)
    p0.font.bold = True
    p0.font.color.rgb = IBM_TEAL
    p0.space_after = Pt(14)

    p1 = tf1.add_paragraph()
    p1.text = "Generative & Agentic AI Consultant Application"
    p1.font.size = Pt(36)
    p1.font.bold = True
    p1.font.color.rgb = TEXT_WHITE
    p1.space_after = Pt(14)

    p2 = tf1.add_paragraph()
    p2.text = "Demonstrating Client-First Value, watsonx & Partner Ecosystem Alignment,\nAugmented Work Products, and Trustworthy AI Governance"
    p2.font.size = Pt(16)
    p2.font.color.rgb = TEXT_LIGHT

    # Metadata Grid
    meta_items = [
        ("Candidate / Role", "Senior GenAI Consultant / BA"),
        ("Target Client", "State Insurance Group (SIG)"),
        ("Badge Target", "Experienced Consultant Level"),
        ("Presentation Time", "7-10 Minutes (~8.5 Mins)")
    ]
    for idx, (label, val) in enumerate(meta_items):
        col_left = 1.0 + (idx * 2.9)
        add_card(slide1, col_left, 5.2, 2.7, 1.3, label, paragraph_text=val, border_color=IBM_BLUE)

    set_speaker_notes(slide1, "Welcome evaluator team. Today I am presenting my Experienced GenAI Consultant Stand & Deliver deck, covering our State Insurance Group engagement.")

    # -------------------------------------------------------------
    # SLIDE 2: CLIENT NEED & BUSINESS PROBLEM
    # -------------------------------------------------------------
    slide2 = prs.slides.add_slide(blank_layout)
    set_background(slide2)
    add_header(slide2, "1. Client Need & Business Problem", "Rubric Area 1")

    add_card(slide2, 0.6, 1.4, 5.9, 5.3, "Client Context & Objectives", 
             paragraph_text="Client: State Insurance Group (Claims Modernization Program)",
             items=[
                 "Legacy Debt: Claims processing platform relied on 15+ year old legacy rules with zero updated technical specs.",
                 "Strategic Goal: Modernize claims intake and policy validation while ensuring 100% HIPAA and state compliance.",
                 "Business Need: Accelerate discovery, requirement gathering, and sprint delivery speed without headcount expansion."
             ])

    add_card(slide2, 6.8, 1.4, 5.9, 5.3, "Consultant & BA Pain Points", 
             items=[
                 "Requirement Bottlenecks: Business Analysts spent 60%+ of sprint cycles manually drafting user stories and test scenarios.",
                 "Knowledge Gaps: Inconsistent business logic documentation led to frequent rework during development.",
                 "Stakeholder Alignment: Slow translation between non-technical claims adjusters and software engineering teams."
             ])

    set_speaker_notes(slide2, "My role on the State Insurance Group engagement was lead GenAI Consultant and Business Analyst. Our client faced a massive documentation backlog and slow claims processing workflows.")

    # -------------------------------------------------------------
    # SLIDE 3: RELEVANT OFFERING & AI ECOSYSTEM
    # -------------------------------------------------------------
    slide3 = prs.slides.add_slide(blank_layout)
    set_background(slide3)
    add_header(slide3, "2. Relevant Offering & AI Ecosystem", "Rubric Area 2")

    add_card(slide3, 0.6, 1.4, 3.8, 5.3, "IBM Consulting Offering", 
             paragraph_text="GenAI Application Modernization:\nEnabling end-to-end SDLC augmentation using IBM Consulting Advantage assets to accelerate discovery, design, and story writing.")

    add_card(slide3, 4.75, 1.4, 3.8, 5.3, "IBM watsonx Portfolio", 
             items=[
                 "watsonx.ai: Granite 3.0 models for enterprise prompt engineering & document extraction.",
                 "watsonx Orchestrate: Multi-agent workflows for automated approval routing.",
                 "watsonx.governance: Risk tracking & auditability."
             ])

    add_card(slide3, 8.9, 1.4, 3.8, 5.3, "Strategic Partner Stack", 
             paragraph_text="Microsoft Azure OpenAI & Copilot:\nIntegrated into developer IDEs and Microsoft 365, complementing watsonx for code synthesis and stakeholder documentation summarization.")

    set_speaker_notes(slide3, "We brought a 'Client-first with a point of view' perspective. We combined IBM watsonx.ai Granite models with Azure OpenAI, governed by watsonx.governance.")

    # -------------------------------------------------------------
    # SLIDE 4: WORK PRODUCTS: AUGMENTATION
    # -------------------------------------------------------------
    slide4 = prs.slides.add_slide(blank_layout)
    set_background(slide4)
    add_header(slide4, "3. Work Products: 50%-75% Work Augmentation", "Rubric Area 3")

    add_card(slide4, 0.6, 1.4, 6.8, 5.3, "Augmenting Business Analysis Workflows", 
             paragraph_text="As Consultant/BA, I utilized structured prompt engineering and IBM Assistants to augment 60%+ of core deliverables:",
             items=[
                 "Auto-Generated User Stories: Transformed unstructured claim policy PDFs into Jira-ready User Stories with Gherkin Acceptance Criteria.",
                 "Process Flow Mapping: Synthesized complex business process models into structured BPMN specifications in minutes.",
                 "Synthetic Test Data Generation: Created realistic, anonymized claims datasets for integration testing."
             ])

    # Metric Card Callout
    metric_card = slide4.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(7.7), Inches(1.4), Inches(5.0), Inches(5.3))
    metric_card.fill.solid()
    metric_card.fill.fore_color.rgb = CARD_BG
    metric_card.line.color.rgb = IBM_BLUE
    metric_card.line.width = Pt(2)

    mtf = metric_card.text_frame
    mtf.word_wrap = True
    mp1 = mtf.paragraphs[0]
    mp1.alignment = PP_ALIGN.CENTER
    mp1.text = "65%"
    mp1.font.size = Pt(64)
    mp1.font.bold = True
    mp1.font.color.rgb = IBM_TEAL

    mp2 = mtf.add_paragraph()
    mp2.alignment = PP_ALIGN.CENTER
    mp2.text = "Reduction in Requirement Drafting Time\n"
    mp2.font.size = Pt(14)
    mp2.font.color.rgb = TEXT_WHITE

    mp3 = mtf.add_paragraph()
    mp3.text = "Deliverable Produced: Claims Processing Agile Backlog & Functional Requirements Specification (FRS) created in 3 days vs. 3 weeks."
    mp3.font.size = Pt(13)
    mp3.font.color.rgb = TEXT_LIGHT

    set_speaker_notes(slide4, "I built a specialized Prompt Card within IBM Consulting Advantage that ingested legacy policy documents and automatically generated structured Agile user stories.")

    # -------------------------------------------------------------
    # SLIDE 5: AGENTIC AI INTEGRATION
    # -------------------------------------------------------------
    slide5 = prs.slides.add_slide(blank_layout)
    set_background(slide5)
    add_header(slide5, "3. Agentic AI Integration & Architecture", "Rubric Area 3")

    add_card(slide5, 0.6, 1.4, 3.8, 5.3, "Agent 1: Intake & Extraction", 
             paragraph_text="watsonx.ai Agent:\nExtracts key claim entities (policy number, claim amount, incident date) from unstructured adjuster notes with 98% accuracy.")

    add_card(slide5, 4.75, 1.4, 3.8, 5.3, "Agent 2: Policy Decisioning", 
             paragraph_text="watsonx Orchestrate Agent:\nCross-references claim details against active policy terms and evaluates fraud risk thresholds automatically.")

    add_card(slide5, 8.9, 1.4, 3.8, 5.3, "Human-in-the-Loop (HITL)", 
             paragraph_text="Escalation Gateway:\nHigh-value or flagged claims are seamlessly passed to human claims adjusters with pre-generated reasoning summaries and recommendation flags.")

    set_speaker_notes(slide5, "Beyond basic GenAI, we deployed an Agentic AI workflow using watsonx Orchestrate. Autonomous agents collaborated to perform claims validation and entity extraction.")

    # -------------------------------------------------------------
    # SLIDE 6: IBM METHODS, TOOLS & BEST PRACTICES
    # -------------------------------------------------------------
    slide6 = prs.slides.add_slide(blank_layout)
    set_background(slide6)
    add_header(slide6, "4. IBM Methods, Tools & Best Practices", "Rubric Area 4")

    steps = [
        ("METHOD 01", "IBM Garage Co-Creation", "Executed 2-week MVP innovation sprints with insurance stakeholders, rapid prototyping prompt concepts and user feedback loops."),
        ("METHOD 02", "IBM Consulting Advantage", "Utilized standard IBM AI Assistants, Prompt Libraries, and Delivery Method cards to maintain consistent delivery governance."),
        ("METHOD 03", "IBM Core Delivery Method", "Embedded AI checks directly into Definition of Ready (DoR) and Definition of Done (DoD) agile quality gate criteria.")
    ]

    for idx, (num, title, desc) in enumerate(steps):
        left = 0.6 + (idx * 4.15)
        add_card(slide6, left, 1.4, 3.8, 5.3, f"{num}\n{title}", paragraph_text=desc, border_color=IBM_TEAL)

    set_speaker_notes(slide6, "We followed IBM Garage principles to co-create solution prototypes directly with claims adjusters, utilizing IBM Consulting Advantage prompt cards.")

    # -------------------------------------------------------------
    # SLIDE 7: TRUSTWORTHY AI, RISKS & ETHICS
    # -------------------------------------------------------------
    slide7 = prs.slides.add_slide(blank_layout)
    set_background(slide7)
    add_header(slide7, "5. Trustworthy AI, Risks & Ethics", "Rubric Area 5")

    add_card(slide7, 0.6, 1.4, 5.9, 5.3, "Identified AI Risks in Insurance", 
             items=[
                 "Data Privacy & PII Leakage: Exposure of policyholder social security numbers or private health information (PHI).",
                 "Algorithmic Bias: Risk of model demographic bias during automated claims decisioning.",
                 "Hallucination & Drift: Generative models outputting incorrect policy terms or invalid clause references."
             ])

    add_card(slide7, 6.8, 1.4, 5.9, 5.3, "IBM Governance & Ethical Controls", 
             items=[
                 "watsonx.governance Tracking: Monitored model lineage, prompt drift, and fairness metrics in real-time.",
                 "Data Masking & Zero Retention: Implemented automated PII redaction filters before LLM ingestion.",
                 "5 Pillars of Trustworthy AI: Enforced Explainability, Fairness, Robustness, Transparency, and Privacy across all agents."
             ])

    set_speaker_notes(slide7, "Financial services demand uncompromising ethics. We addressed data privacy using automated PII masking in watsonx.ai and tracked explainability in watsonx.governance.")

    # -------------------------------------------------------------
    # SLIDE 8: MEASURABLE OUTCOMES & BUSINESS VALUE
    # -------------------------------------------------------------
    slide8 = prs.slides.add_slide(blank_layout)
    set_background(slide8)
    add_header(slide8, "6. Measurable Outcomes & Business Value", "Value Impact")

    metrics = [
        ("45%", "Faster Backlog Refinement"),
        ("30%", "Increase in Sprint Velocity"),
        ("99.2%", "Compliance & Audit Pass Rate")
    ]
    for idx, (num, label) in enumerate(metrics):
        left = 0.6 + (idx * 4.15)
        mcard = slide8.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(left), Inches(1.4), Inches(3.8), Inches(1.8))
        mcard.fill.solid()
        mcard.fill.fore_color.rgb = CARD_BG
        mcard.line.color.rgb = IBM_BLUE
        
        tf = mcard.text_frame
        p1 = tf.paragraphs[0]
        p1.alignment = PP_ALIGN.CENTER
        p1.text = num
        p1.font.size = Pt(36)
        p1.font.bold = True
        p1.font.color.rgb = IBM_TEAL
        
        p2 = tf.add_paragraph()
        p2.alignment = PP_ALIGN.CENTER
        p2.text = label
        p2.font.size = Pt(13)
        p2.font.color.rgb = TEXT_WHITE

    add_card(slide8, 0.6, 3.5, 12.133, 3.2, "Delivered Client Value to State Insurance Group", 
             paragraph_text="Accelerated the Claims Modernization roadmap by 2.5 months while maintaining zero high-severity compliance or security findings. Reallocated 200+ BA and developer hours per sprint toward strategic business architecture and complex claims handling.")

    set_speaker_notes(slide8, "The business impact was undeniable. Sprint velocity surged by 30%, requirement discovery time plummeted by 45%, and the client achieved a 99.2% compliance pass rate.")

    # -------------------------------------------------------------
    # SLIDE 9: REUSABLE ASSETS & EMINENCE
    # -------------------------------------------------------------
    slide9 = prs.slides.add_slide(blank_layout)
    set_background(slide9)
    add_header(slide9, "7. Reusable Assets & Community Giveback", "Eminence")

    add_card(slide9, 0.6, 1.4, 5.9, 5.3, "Created Reusable Assets", 
             items=[
                 "Insurance GenAI Prompt Library: 25+ vetted prompt templates for business logic extraction and test scenario generation.",
                 "watsonx Orchestrate Agent Template: Reusable multi-agent workflow blueprint for insurance approval chains.",
                 "BA GenAI Onboarding Guide: Step-by-step playbook published on IBM Consulting Advantage asset hub."
             ])

    add_card(slide9, 6.8, 1.4, 5.9, 5.3, "Mentorship & Practice Eminence", 
             items=[
                 "Practice Community Tech Talk: Led a global IBM GenAI Community of Practice session with 180+ attendees.",
                 "Consultant Mentorship: Upskilled 6 junior Business Analysts and Consultants on prompt engineering techniques.",
                 "Thought Leadership: Co-authored internal IBM case study on Agentic Workflows in Financial Services."
             ])

    set_speaker_notes(slide9, "To drive eminence and give back to IBM, I packaged our prompt templates and agentic workflow blueprints into reusable assets on IBM Consulting Advantage.")

    # -------------------------------------------------------------
    # SLIDE 10: CONCLUSION & SUMMARY
    # -------------------------------------------------------------
    slide10 = prs.slides.add_slide(blank_layout)
    set_background(slide10)
    add_header(slide10, "Conclusion & Badge Readiness", "Wrap-Up")

    c_box = slide10.shapes.add_textbox(Inches(1.5), Inches(1.5), Inches(10.333), Inches(1.8))
    ctf = c_box.text_frame
    ctf.word_wrap = True
    
    cp1 = ctf.paragraphs[0]
    cp1.alignment = PP_ALIGN.CENTER
    cp1.text = "Delivering Superior Value with IBM GenAI & Agentic AI"
    cp1.font.size = Pt(24)
    cp1.font.bold = True
    cp1.font.color.rgb = IBM_TEAL
    cp1.space_after = Pt(10)

    cp2 = ctf.add_paragraph()
    cp2.alignment = PP_ALIGN.CENTER
    cp2.text = "Demonstrated full alignment across all 5 IBM Experienced Consultant evaluation criteria through proven client impact, watsonx & partner ecosystem mastery, robust Trustworthy AI governance, and reusable asset creation."
    cp2.font.size = Pt(14)
    cp2.font.color.rgb = TEXT_LIGHT

    summary_cards = [
        ("Client First Value", "30% Sprint Acceleration"),
        ("Trustworthy AI", "watsonx.governance Control"),
        ("IBM Eminence", "Assets & Mentorship")
    ]
    for idx, (title, sub) in enumerate(summary_cards):
        left = 1.0 + (idx * 3.9)
        add_card(slide10, left, 3.6, 3.5, 3.0, title, paragraph_text=sub, border_color=IBM_BLUE)

    set_speaker_notes(slide10, "In conclusion, this project exemplifies IBM's value proposition—combining cutting-edge watsonx technology with deep domain consulting expertise. Thank you.")

    # Save presentation
    output_path = "IBM_GenAI_Consultant_Presentation.pptx"
    prs.save(output_path)
    print(f"Presentation saved successfully to {os.path.abspath(output_path)}")

if __name__ == "__main__":
    create_deck()